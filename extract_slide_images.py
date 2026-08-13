import pptx
import os

prs = pptx.Presentation(r'C:\Users\felix\Downloads\ClassSync_CTEACH_Felix_Alabi.pptx')
out_dir = r'C:\Users\felix\.gemini\antigravity\scratch\ClassSync-Product-Launch\extracted_slides'
os.makedirs(out_dir, exist_ok=True)

for idx, slide in enumerate(prs.slides):
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.ext
            filename = f"slide_{idx+1}_pic_{shape_idx+1}.{ext}"
            filepath = os.path.join(out_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(image.blob)
            print(f"Saved {filepath}")
