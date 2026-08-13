import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = pptx.Presentation(r'C:\Users\felix\Downloads\ClassSync_CTEACH_Felix_Alabi.pptx')

def inspect_shape(shape, depth=0):
    indent = "  " * depth
    print(f"{indent}- Shape: '{shape.name}', Type: {shape.shape_type}")
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if txt:
            print(f"{indent}  Text: {txt}")
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            inspect_shape(s, depth + 1)

for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx+1} ====================")
    for shape in slide.shapes:
        inspect_shape(shape)
