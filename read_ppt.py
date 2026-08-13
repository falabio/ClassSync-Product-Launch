import pptx

prs = pptx.Presentation(r'C:\Users\felix\Downloads\ClassSync_CTEACH_Felix_Alabi.pptx')
print(f"Total slides: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx+1} ====================")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"[Shape Text]: {text}")
        if shape.has_table:
            print("[Table]:")
            for row in shape.table.rows:
                cell_texts = [cell.text.strip() for cell in row.cells]
                print("  | " + " | ".join(cell_texts) + " |")
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"[Speaker Notes]: {notes}")
